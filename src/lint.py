"""Geometry linter: flags text that overflows its box and shapes off-slide."""

import sys
from pptx import Presentation
from pptx.util import Emu
import deck as D

EMU = 914400.0


def text_of(tf):
    return "\n".join(p.text for p in tf.paragraphs)


def lint(path, tol=0.06):
    prs = Presentation(path)
    sw = prs.slide_width / EMU
    sh = prs.slide_height / EMU
    issues = []
    logos = {}
    for idx, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            try:
                x, y = shp.left / EMU, shp.top / EMU
                w, h = shp.width / EMU, shp.height / EMU
            except TypeError:
                continue
            deco = (shp.name or "").startswith("deco-")
            if not deco and (x < -0.02 or y < -0.02
                             or x + w > sw + 0.02 or y + h > sh + 0.02):
                issues.append((idx, "off-slide", f"{shp.shape_type} "
                                                 f"@({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}"))
            if (shp.name or "") == "college-logo":
                logos[idx] = logos.get(idx, 0) + 1
            if not shp.has_text_frame:
                continue
            tf = shp.text_frame
            txt = text_of(tf)
            if not txt.strip():
                continue
            ml = (tf.margin_left or 0) / EMU
            mr = (tf.margin_right or 0) / EMU
            mt = (tf.margin_top or 0) / EMU
            mb = (tf.margin_bottom or 0) / EMU
            iw = max(w - ml - mr, 0.3)
            ih = max(h - mt - mb, 0.1)
            total = 0.0
            for p in tf.paragraphs:
                ptxt = p.text
                if not ptxt:
                    continue
                runs = p.runs
                size = None
                bold = False
                family = "sans"
                for r in runs:
                    if r.font.size:
                        size = r.font.size.pt
                        bold = bool(r.font.bold)
                        if (r.font.name or "") == D.SERIF:
                            family = "serif"
                        break
                size = size or 18.0
                ls = p.line_spacing or 1.2
                if not isinstance(ls, float):
                    ls = 1.2
                sb = (p.space_before.pt if p.space_before else 0) / 72.0
                sa = (p.space_after.pt if p.space_after else 0) / 72.0
                lines = D.fit(ptxt, size, iw, bold=bold, family=family)
                total += lines * D.line_height_in(size, ls, family) + sb + sa
            if total > ih + tol:
                issues.append((idx, "overflow",
                               f"needs {total:.2f}in, box {ih:.2f}in :: "
                               f"{txt[:58]!r}"))
    for idx, count in sorted(logos.items()):
        if count > 1:
            issues.append((idx, "dup-logo", f"{count} logo images on one slide"))
    return issues, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    total = 0
    for path in sys.argv[1:]:
        issues, n = lint(path)
        print(f"\n=== {path.split('/')[-1]}  ({n} slides) ===")
        if not issues:
            print("  clean")
        for idx, kind, msg in issues:
            print(f"  slide {idx:>2}  {kind:<9} {msg}")
        total += len(issues)
    print(f"\nTOTAL ISSUES: {total}")
