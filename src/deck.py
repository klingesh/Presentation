"""
deck.py — corporate presentation design system (ISSM report theme).

Styled to match the ISSM Business School report template: crimson and charcoal
on white, Sorts Mill Goudy for headings and figures, Quattrocento Sans for
body copy, numbered section badges, panel-and-card anatomy and small captions.

The public API (Deck + its layout methods) is deliberately unchanged, so the
per-unit content files drive this theme without edits.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# --------------------------------------------------------------------------
# Palette — crimson primary, charcoal secondary, warm greys
# --------------------------------------------------------------------------
CRIMSON      = RGBColor(0x8B, 0x00, 0x00)
CRIMSON_DK   = RGBColor(0x63, 0x00, 0x00)
CRIMSON_MID  = RGBColor(0xA8, 0x2B, 0x2B)
CRIMSON_TINT = RGBColor(0xF8, 0xEE, 0xEE)

CHARCOAL     = RGBColor(0x1F, 0x29, 0x37)
CHARCOAL_DK  = RGBColor(0x0F, 0x17, 0x22)
CHARCOAL_TINT = RGBColor(0xEE, 0xF0, 0xF3)

STEEL        = RGBColor(0x4B, 0x55, 0x63)
STEEL_TINT   = RGBColor(0xF1, 0xF2, 0xF4)

WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
PANEL        = RGBColor(0xF9, 0xFA, 0xFB)
BORDER       = RGBColor(0xE5, 0xE7, 0xEB)
INK          = RGBColor(0x1F, 0x29, 0x37)
MUTED        = RGBColor(0x6B, 0x72, 0x80)
MUTED_LT     = RGBColor(0x9C, 0xA3, 0xAF)

# text on dark backgrounds
ON_DARK      = RGBColor(0xC7, 0xCE, 0xD8)
ON_CRIMSON   = RGBColor(0xF2, 0xD9, 0xD9)
ON_CRIMSON_LT = RGBColor(0xE3, 0xB4, 0xB4)

# Legacy names kept so the unit content files need no changes. The rainbow
# palette collapses onto the report's restrained crimson / charcoal / steel.
NAVY = INK
BLUE = CRIMSON
RED = CRIMSON
AMBER = CRIMSON
TEAL = CHARCOAL
GREEN = CHARCOAL
VIOLET = STEEL
BLUE_LIGHT = CRIMSON_TINT
TEAL_LIGHT = CHARCOAL_TINT
AMBER_LIGHT = CRIMSON_TINT
GREEN_LIGHT = CHARCOAL_TINT
VIOLET_LIGHT = STEEL_TINT
RED_LIGHT = CRIMSON_TINT
NAVY_SOFT = STEEL

ACCENTS = [CRIMSON, CHARCOAL, STEEL, CRIMSON, CHARCOAL, STEEL]
ACCENT_TINTS = {
    str(CRIMSON): CRIMSON_TINT,
    str(CHARCOAL): CHARCOAL_TINT,
    str(STEEL): STEEL_TINT,
}


def tint(color):
    return ACCENT_TINTS.get(str(color), CRIMSON_TINT)


# --------------------------------------------------------------------------
# Typography
# --------------------------------------------------------------------------
SERIF = "Sorts Mill Goudy"      # headings, titles, big figures (no bold cut)
SANS = "Quattrocento Sans"      # body copy, labels, captions
FONT = SANS                     # back-compat alias

FONT_DIR = "/usr/share/fonts/issm"
FONT_FILES = {
    ("sans", False): f"{FONT_DIR}/QuattrocentoSans-Regular.ttf",
    ("sans", True): f"{FONT_DIR}/QuattrocentoSans-Bold.ttf",
    # Sorts Mill Goudy ships regular + italic only; bold is synthesised
    ("serif", False): f"{FONT_DIR}/SortsMillGoudy-Regular.ttf",
    ("serif", True): f"{FONT_DIR}/SortsMillGoudy-Regular.ttf",
}
# PowerPoint percentage line spacing scales the font's own line height
LINE_FACTORS = {"sans": 1.108, "serif": 1.438}
TYPEFACE = {"sans": SANS, "serif": SERIF}

_PIL_CACHE = {}
_SCALE = 8


def _pil_font(size, bold, family):
    key = (round(float(size), 2), bool(bold), family)
    if key not in _PIL_CACHE:
        from PIL import ImageFont
        _PIL_CACHE[key] = ImageFont.truetype(
            FONT_FILES[(family, bool(bold))], max(1, int(round(size * _SCALE))))
    return _PIL_CACHE[key]


def text_width_in(text, size, bold=False, family="sans"):
    if not text:
        return 0.0
    try:
        f = _pil_font(size, bold, family)
        w = (f.getlength(text) / _SCALE) / 72.0
        if family == "serif" and bold:
            w *= 1.03           # allow for synthesised bold
        return w
    except Exception:
        return len(text) * size * 0.5 / 72.0


def fit(text, size, width_in, bold=False, family="sans"):
    """Number of wrapped lines the text needs inside width_in."""
    total = 0
    for hard in str(text).replace("\v", "\n").replace("\x0b", "\n").split("\n"):
        if not hard.strip():
            total += 1
            continue
        lines, cur = 1, ""
        for w in hard.split():
            trial = (cur + " " + w).strip()
            if cur and text_width_in(trial, size, bold, family) > width_in:
                lines += 1
                cur = w
            else:
                cur = trial
        total += lines
    return max(total, 1)


def line_height_in(size, line=1.25, family="sans"):
    return size * LINE_FACTORS[family] * line / 72.0


def text_height_in(text, size, width_in, bold=False, line=1.25, family="sans"):
    return fit(text, size, width_in, bold, family) * line_height_in(size, line,
                                                                   family)


def autosize(text, width_in, height_in, size, bold=False, line=1.25,
             min_size=8.0, step=0.5, family="sans"):
    s = float(size)
    while s > min_size and text_height_in(text, s, width_in, bold, line,
                                         family) > height_in:
        s -= step
    return s


# --------------------------------------------------------------------------
# Grid (16:9 — 13.333 x 7.5 in)
# --------------------------------------------------------------------------
SW, SH = 13.333, 7.5
ML, MR = 0.62, 0.62
CW = SW - ML - MR                # 12.093
BODY_INDENT = 0.30               # body sits right of the vertical crimson rule
BODY_X = ML + BODY_INDENT
BODY_W = CW - BODY_INDENT
BODY_BOTTOM = 6.66
FOOTER_Y = 6.94


def _i(v):
    return Inches(v) if not isinstance(v, Emu) else v


# --------------------------------------------------------------------------
# Shape helpers — flat, bordered, square-cornered (report style)
# --------------------------------------------------------------------------
def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75,
         shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=False):
    s = slide.shapes.add_shape(shape, _i(x), _i(y), _i(w), _i(h))
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    if shadow:
        soft_shadow(s)
    s.text_frame.word_wrap = True
    return s


def card(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=None, shadow=False):
    """Flat bordered panel — the workhorse of this theme."""
    return rect(slide, x, y, w, h, fill=fill, line=line, lw=0.75)


def circle(slide, cx, cy, d, fill=CRIMSON, line=None):
    return rect(slide, cx - d / 2, cy - d / 2, d, d, fill=fill, line=line,
                shape=MSO_SHAPE.OVAL)


def hline(slide, x, y, w, color=BORDER, weight=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), _i(w),
                               Pt(weight))
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def vline(slide, x, y, h, color=CRIMSON, weight=1.6):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), Pt(weight),
                               _i(h))
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def marker(slide, x, y, size=0.085, color=CRIMSON):
    """Small square bullet/heading marker."""
    return rect(slide, x, y, size, size, fill=color)


def soft_shadow(shape, blur=14, dist=2, alpha=10000):
    spPr = shape._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    spPr.append(parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/'
        f'2006/main"><a:outerShdw blurRad="{blur * 12700}" '
        f'dist="{dist * 12700}" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="1F2937"><a:alpha val="{alpha}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'))


def gradient_bg(slide, c1, c2, angle=45.0):
    s = rect(slide, 0, 0, SW, SH, fill=c1)
    s.fill.gradient()
    stops = s.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    s.fill.gradient_angle = angle
    s.line.fill.background()
    return s


def ghost(slide, x, y, w, h, color="FFFFFF", alpha=6000,
          shape=MSO_SHAPE.OVAL):
    s = slide.shapes.add_shape(shape, _i(x), _i(y), _i(w), _i(h))
    s.name = "deco-bleed"
    s.shadow.inherit = False
    s.line.fill.background()
    spPr = s._element.spPr
    for e in spPr.findall(qn("a:solidFill")):
        spPr.remove(e)
    geom = spPr.find(qn("a:prstGeom"))
    geom.addnext(parse_xml(
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/'
        f'2006/main"><a:srgbClr val="{color}"><a:alpha val="{alpha}"/>'
        '</a:srgbClr></a:solidFill>'))
    return s


# --------------------------------------------------------------------------
# Text helpers
# --------------------------------------------------------------------------
def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def style_run(run, size=11, bold=False, color=INK, family="sans", italic=False,
              spacing=None, caps=False):
    f = run.font
    f.name = TYPEFACE[family]
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    if spacing:
        rPr.set("spc", str(int(spacing * 100)))
    if caps:
        rPr.set("cap", "all")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            rPr.append(parse_xml(
                f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/'
                f'2006/main" typeface="{TYPEFACE[family]}"/>'))
        else:
            el.set("typeface", TYPEFACE[family])
    return run


def para(tf, text, size=11, bold=False, color=INK, family="sans",
         align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=1.25,
         italic=False, spacing=None, caps=False, first=False, indent=0.0):
    p = tf.paragraphs[0] if (first or not getattr(tf, "_used", False)) \
        else tf.add_paragraph()
    tf._used = True
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if indent:
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(indent * 914400)))
        pPr.set("indent", str(-int(0.18 * 914400)))
    for j, chunk in enumerate(str(text).split("\n")):
        if j:
            p._p.append(parse_xml(
                '<a:br xmlns:a="http://schemas.openxmlformats.org/drawingml/'
                '2006/main"/>'))
        r = p.add_run()
        r.text = chunk
        style_run(r, size=size, bold=bold, color=color, family=family,
                  italic=italic, spacing=spacing, caps=caps)
    return p


def rich(tf, parts, size=11, color=INK, family="sans", align=PP_ALIGN.LEFT,
         line=1.25, space_before=0, space_after=0, first=False):
    p = tf.paragraphs[0] if (first or not getattr(tf, "_used", False)) \
        else tf.add_paragraph()
    tf._used = True
    p.alignment = align
    p.line_spacing = line
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    for text, ov in parts:
        r = p.add_run()
        r.text = text
        style_run(r, size=ov.get("size", size), bold=ov.get("bold", False),
                  color=ov.get("color", color), family=ov.get("family", family),
                  italic=ov.get("italic", False), spacing=ov.get("spacing"),
                  caps=ov.get("caps", False))
    return p


def eyebrow(slide, x, y, text, color=MUTED, size=8.0, w=6.0):
    """Tiny letter-spaced uppercase label."""
    tf = textbox(slide, x, y, w, 0.2)
    para(tf, text.upper(), size=size, color=color, spacing=1.3, first=True,
         line=1.0)
    return tf


def caption(slide, x, y, text, w=5.4, align=PP_ALIGN.LEFT, color=MUTED_LT):
    tf = textbox(slide, x, y, w, 0.22)
    para(tf, text, size=7.5, color=color, align=align, line=1.1, first=True)
    return tf


def num_badge(slide, x, y, n, size=0.34, fill=CRIMSON, fsize=13,
              family="serif", color=WHITE, rounded=False):
    rect(slide, x, y, size, size, fill=fill,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
         radius=0.18 if rounded else None)
    tf = textbox(slide, x, y + size * 0.5 - fsize * 0.011, size, size)
    para(tf, str(n), size=fsize, color=color, family=family,
         align=PP_ALIGN.CENTER, line=1.0, first=True)


def chip(slide, x, y, text, fill=CRIMSON, color=WHITE, size=8.0, h=0.28,
         pad=0.34, bold=False, family="sans", line=None):
    w = text_width_in(text.upper(), size, bold, family) + \
        (1.3 / 72.0) * max(len(text) - 1, 0) + pad
    box = rect(slide, x, y, w, h, fill=fill, line=line)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = 0
    para(tf, text.upper(), size=size, color=color, family=family, bold=bold,
         align=PP_ALIGN.CENTER, spacing=1.3, line=1.0, first=True)
    return w


def chip_width(text, size=8.0, bold=False, spacing=1.3, pad=0.34,
               family="sans"):
    return (text_width_in(text.upper(), size, bold, family)
            + (spacing / 72.0) * max(len(text) - 1, 0) + pad)


# --------------------------------------------------------------------------
# Logo
# --------------------------------------------------------------------------
EXTS = ("png", "PNG", "jpg", "JPG", "jpeg", "JPEG", "webp", "gif")


def _assets_dir(assets_dir=None):
    import os
    if assets_dir:
        return assets_dir
    return os.path.join(os.path.dirname(os.path.abspath(__file__)), "..",
                        "assets")


def find_logo(prefer="logo-full", assets_dir=None, fallback=True):
    import glob
    import os
    d = _assets_dir(assets_dir)
    for ext in EXTS:
        hit = os.path.join(d, f"{prefer}.{ext}")
        if os.path.exists(hit):
            return hit
    if not fallback:
        return None
    for ext in EXTS:
        hits = sorted(g for g in glob.glob(os.path.join(d, f"*.{ext}"))
                      if "source" not in os.path.basename(g).lower())
        if hits:
            return hits[0]
    return None


def logo_size(path, h, max_w=3.0):
    if not path:
        return 0.0, 0.0
    try:
        from PIL import Image
        pw, ph = Image.open(path).size
        w = h * (pw / float(ph))
    except Exception:
        w = h
    if w > max_w:
        h = h * (max_w / w)
        w = max_w
    return w, h


def place_logo(slide, path, x, y, h, on_dark=False, pad=0.13, max_w=3.0,
               right_edge=None):
    if not path:
        return 0.0
    w, h = logo_size(path, h, max_w)
    if right_edge is not None:
        x = right_edge - w
    if on_dark:
        plate = rect(slide, x - pad, y - pad, w + pad * 2, h + pad * 2,
                     fill=WHITE, line=None)
        plate.name = "logo-plate"
    pic = slide.shapes.add_picture(path, _i(x), _i(y), height=_i(h))
    pic.name = "college-logo"
    return w


# ==========================================================================
# Deck
# ==========================================================================
class Deck:
    PRESENTER = "Lingesh K"
    REGISTER = "OSI2509030"

    def __init__(self, unit_no, unit_title, course="PGPM · Systems Specialisation",
                 subject="Artificial Intelligence & Generative AI", logo=None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SW)
        self.prs.slide_height = Inches(SH)
        self.unit_no = unit_no
        self.unit_title = unit_title
        self.course = course
        self.subject = subject
        self.logo = logo if logo is not None else find_logo("logo-full")
        self.mark = find_logo("logo-mark", fallback=False) or self.logo
        self.blank = self.prs.slide_layouts[6]
        self._section_n = 0
        self._content_n = 0

        core = self.prs.core_properties
        core.title = f"Unit {unit_no} — {unit_title}"
        core.author = f"{self.PRESENTER} ({self.REGISTER})"
        core.subject = subject
        core.comments = f"{course} · Unit {unit_no} of 5"

    # ---------------- infrastructure ----------------
    def _new(self, footer=True, bg=WHITE):
        s = self.prs.slides.add_slide(self.blank)
        if bg is not None:
            rect(s, 0, 0, SW, SH, fill=bg)
        if footer:
            self._footer(s)
        return s

    def _slide_no(self):
        return len(self.prs.slides._sldIdLst)

    def _footer(self, slide):
        """Small captions, in the report's style, plus the logo top right."""
        caption(slide, ML, FOOTER_Y,
                f"{self.PRESENTER}  ·  Register No. {self.REGISTER}", w=5.0)
        tf = textbox(slide, ML + CW - 3.0, FOOTER_Y, 3.0, 0.22)
        para(tf, f"{self._slide_no():02d}", size=8.5, color=CRIMSON,
             family="serif", align=PP_ALIGN.RIGHT, line=1.1, first=True)
        if self.logo:
            place_logo(slide, self.logo, 0, 0.40, 0.46, on_dark=False,
                       max_w=1.25, right_edge=ML + CW)

    def _head(self, slide, kicker, title, sub=None, accent=CRIMSON,
              numbered=True):
        """Report-style header: number badge, eyebrow, serif title, rule."""
        y = 0.52
        tx = ML
        if numbered:
            self._content_n += 1
            num_badge(slide, ML, y + 0.03, self._content_n, size=0.34,
                      fill=accent, fsize=13.5, family="serif")
            tx = ML + 0.50

        tw = CW - (tx - ML) - 1.55           # keep clear of the logo
        if kicker:
            eyebrow(slide, tx, y, kicker, color=MUTED, size=8.0, w=tw)
            ty = y + 0.20
        else:
            ty = y + 0.02

        tsize = 25.0
        while tsize > 17 and fit(title, tsize, tw, family="serif") > 1:
            tsize -= 1.0
        th = text_height_in(title, tsize, tw, False, 1.0, "serif")
        tf = textbox(slide, tx, ty, tw, th + 0.12)
        para(tf, title, size=tsize, color=INK, family="serif", line=1.0,
             first=True)

        rule_y = max(ty + th + 0.14, y + 0.44)
        hline(slide, ML, rule_y, CW, accent, 1.25)
        body_top = rule_y + 0.30

        if sub:
            sh = text_height_in(sub, 10.5, CW * 0.82, False, 1.3)
            tf = textbox(slide, ML, rule_y + 0.16, CW * 0.82, sh + 0.1)
            para(tf, sub, size=10.5, color=MUTED, line=1.3, first=True)
            body_top = rule_y + 0.16 + sh + 0.24

        vline(slide, ML, body_top - 0.04, BODY_BOTTOM - body_top + 0.04,
              accent, 1.5)
        return body_top

    def save(self, path):
        self.prs.save(path)
        return path

    # ---------------- shared building blocks ----------------
    def _panel_head(self, slide, x, y, text, w, accent=CRIMSON, size=12.5,
                    family="serif"):
        """Marker + heading, as used on every panel in the reference deck."""
        marker(slide, x, y + 0.055, 0.085, accent)
        tf = textbox(slide, x + 0.19, y - 0.03, w - 0.19,
                     line_height_in(size, 1.1, family) + 0.1)
        para(tf, text, size=size, color=INK, family=family, line=1.1,
             first=True)
        return text_height_in(text, size, w - 0.19, False, 1.1, family) + 0.1

    def _band(self, slide, y, text, label=None, fill=CRIMSON, h=0.62,
              color=WHITE, label_color=ON_CRIMSON_LT):
        """Full-width crimson callout strip."""
        rect(slide, ML, y, CW, h, fill=fill)
        tf = textbox(slide, ML + 0.3, y, CW - 0.6, h, anchor=MSO_ANCHOR.MIDDLE)
        parts = []
        if label:
            parts.append((label.upper() + "   ", {"size": 8.0,
                                                  "color": label_color,
                                                  "spacing": 1.3}))
        parts.append((text, {"size": 11.5, "color": color}))
        rich(tf, parts, first=True)

    def _dot_rows(self, slide, x, y, w, rows, accent=CRIMSON, avail=None,
                  size=11.0, gap=0.0, boxed=False):
        """Bulleted rows with small crimson dots (optionally in pill boxes)."""
        n = max(len(rows), 1)
        rowh = (avail / n) if avail else 0.34
        psize = min([autosize(r, w - 0.30, rowh - 0.06, size, False, 1.25,
                              min_size=8.5) for r in rows] or [size])
        for r in rows:
            if boxed:
                rect(slide, x, y + 0.01, w, rowh - 0.08, fill=WHITE,
                     line=BORDER)
            circle(slide, x + (0.20 if boxed else 0.055),
                   y + (rowh - 0.08) / 2 if boxed else y + 0.13, 0.075, accent)
            tfp = textbox(slide, x + (0.36 if boxed else 0.22),
                          y + (rowh - 0.08) / 2 - line_height_in(psize, 1.25) / 2
                          if boxed else y - 0.005,
                          w - (0.52 if boxed else 0.30), rowh)
            para(tfp, r, size=psize, color=INK, line=1.25, first=True)
            y += rowh + gap
        return y

    # ================= LAYOUTS =================

    def title_slide(self, title, subtitle, meta_lines=None, chips=None):
        s = self._new(footer=False, bg=None)
        gradient_bg(s, CRIMSON, CRIMSON_DK, 315.0)
        ghost(s, 9.2, -1.9, 6.4, 6.4, "FFFFFF", 3200)
        ghost(s, 11.4, 3.6, 4.6, 4.6, "FFFFFF", 2200)
        ghost(s, -1.5, 4.4, 4.8, 4.8, "000000", 2600)
        if self.logo:
            place_logo(s, self.logo, 0, 0.62, 0.82, on_dark=True, max_w=2.3,
                       right_edge=ML + CW - 0.06)

        x = ML + 0.32
        # eyebrow chip naming the subject
        chip(s, x, 1.66, self.subject, fill=CRIMSON_MID, color=WHITE, size=8.5,
             h=0.32, pad=0.44)

        tw = 9.35
        tsize = 40.0
        while tsize > 28 and fit(title, tsize, tw, family="serif") > 2:
            tsize -= 2.0
        th = text_height_in(title, tsize, tw, False, 1.02, "serif")

        sw = 8.7
        ssize = autosize(subtitle, sw, 1.25, 13.5, False, 1.4, min_size=11)
        subh = text_height_in(subtitle, ssize, sw, False, 1.4)

        chips_h = 0.32 if chips else 0.0
        group = th + 0.28 + subh + (0.40 + chips_h if chips else 0.0)
        top = 2.30 + max(0.0, (5.45 - 2.30 - group) / 2)

        tf = textbox(s, x, top, tw, th + 0.16)
        para(tf, title, size=tsize, color=WHITE, family="serif", line=1.02,
             first=True)

        sy = top + th + 0.28
        tf = textbox(s, x, sy, sw, subh + 0.14)
        para(tf, subtitle, size=ssize, color=ON_CRIMSON, line=1.4, first=True)

        if chips:
            cx = x
            cy = sy + subh + 0.40
            for c in chips:
                w = chip(s, cx, cy, c, fill=None, color=ON_CRIMSON_LT,
                         size=8.0, h=chips_h, pad=0.40, line=CRIMSON_MID)
                cx += w + 0.13

        # presenter panels, in the reference deck's bottom-box style
        boxes = [("Presented by", self.PRESENTER),
                 ("Register No.", self.REGISTER)]
        for m in (meta_lines or []):
            boxes.append(("Course", m))
        bx, by, bh = x, 5.88, 0.72
        for label, value in boxes:
            vw = max(text_width_in(value, 14, False, "serif"),
                     text_width_in(label, 8, False, "sans")) + 0.42
            vw = max(vw, 1.9)
            vline(s, bx, by, bh, CRIMSON_MID, 1.6)
            tf = textbox(s, bx + 0.20, by + 0.06, vw, 0.2)
            para(tf, label.upper(), size=8.0, color=ON_CRIMSON_LT, spacing=1.3,
                 line=1.0, first=True)
            tf = textbox(s, bx + 0.20, by + 0.30, vw, 0.34)
            para(tf, value, size=14, color=WHITE, family="serif", line=1.0,
                 first=True)
            bx += vw + 0.34
        return s

    def agenda_slide(self, items, title="What we will cover",
                     kicker="Agenda", note=None):
        s = self._new()
        top = self._head(s, kicker, title, note)
        n = len(items)
        cols = 2 if n > 4 else 1
        rows = (n + cols - 1) // cols
        gap = 0.32
        cw = (BODY_W - gap) / cols if cols == 2 else BODY_W * 0.8
        avail = BODY_BOTTOM - top
        ch = min(0.92, avail / rows)
        y0 = top + max(0.0, (avail - ch * rows) / 2)
        for i, it in enumerate(items):
            head, desc = (it if isinstance(it, tuple) else (it, None))
            r, c = (i % rows, i // rows) if cols == 2 else (i, 0)
            x = BODY_X + c * (cw + gap)
            y = y0 + r * ch
            num_badge(s, x, y + 0.04, f"{i + 1:02d}", size=0.34, fill=CRIMSON,
                      fsize=11.5, family="serif")
            iw = cw - 0.52
            hsize = autosize(head, iw, ch * (0.46 if desc else 0.9), 12.5, True,
                             1.2, min_size=10.5)
            hh = text_height_in(head, hsize, iw, True, 1.2)
            tf = textbox(s, x + 0.52, y, iw, ch - 0.06)
            para(tf, head, size=hsize, bold=True, color=INK, line=1.2,
                 first=True)
            if desc:
                dsize = autosize(desc, iw, ch - hh - 0.16, 9.5, False, 1.28,
                                 min_size=8.0)
                para(tf, desc, size=dsize, color=MUTED, line=1.28,
                     space_before=2)
        return s

    def section_slide(self, number, title, blurb=None, accent=CRIMSON):
        s = self._new(footer=False, bg=None)
        gradient_bg(s, CHARCOAL, CHARCOAL_DK, 315.0)
        ghost(s, 9.6, -1.4, 5.8, 5.8, "FFFFFF", 2400)
        ghost(s, 11.6, 4.0, 3.8, 3.8, "8B0000", 9000)
        if self.logo:
            place_logo(s, self.logo, 0, 0.62, 0.78, on_dark=True, max_w=2.2,
                       right_edge=ML + CW - 0.06)
        x = ML + 0.32
        # oversized faint section numeral fills the right of the divider
        tf = textbox(s, 7.9, 1.28, 4.9, 4.4)
        para(tf, f"{number:02d}", size=190, color=RGBColor(0x27, 0x31, 0x40),
             family="serif", align=PP_ALIGN.RIGHT, line=1.0, first=True)

        tw = 8.0
        tsize = autosize(title, tw, 2.0, 34, False, 1.04, min_size=24,
                         family="serif")
        th = text_height_in(title, tsize, tw, False, 1.04, "serif")
        bh = text_height_in(blurb, 12.5, 7.4, False, 1.38) if blurb else 0.0
        group = 0.30 + 0.34 + th + 0.24 + (0.30 + bh if blurb else 0.0)
        top = 2.05 + max(0.0, (6.00 - 2.05 - group) / 2)

        chip(s, x, top, f"Section {number:02d}", fill=CRIMSON, color=WHITE,
             size=8.5, h=0.30, pad=0.44)
        ty = top + 0.30 + 0.34
        tf = textbox(s, x, ty, tw, th + 0.16)
        para(tf, title, size=tsize, color=WHITE, family="serif", line=1.04,
             first=True)
        y = ty + th + 0.24
        hline(s, x, y, 0.86, CRIMSON, 1.6)
        if blurb:
            tf = textbox(s, x, y + 0.30, 7.4, bh + 0.12)
            para(tf, blurb, size=12.5, color=ON_DARK, line=1.38, first=True)
        caption(s, x, FOOTER_Y,
                f"{self.PRESENTER}  ·  Register No. {self.REGISTER}", w=5.0,
                color=RGBColor(0x7A, 0x84, 0x93))
        return s

    def bullets_slide(self, kicker, title, bullets, sub=None, lead=None,
                      accent=CRIMSON, two_col=False):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        if lead:
            lsize = autosize(lead, BODY_W - 0.7, 0.44, 11.5, False, 1.25,
                             min_size=9.5)
            rect(s, BODY_X, y, BODY_W, 0.62, fill=accent)
            tf = textbox(s, BODY_X + 0.28, y, BODY_W - 0.56, 0.62,
                         anchor=MSO_ANCHOR.MIDDLE)
            para(tf, lead, size=lsize, color=WHITE, line=1.25, first=True)
            y += 0.62 + 0.30

        cols = 2 if two_col else 1
        rows = (len(bullets) + cols - 1) // cols
        gap = 0.50
        cw = (BODY_W - gap) / cols if cols == 2 else BODY_W * 0.94
        avail = BODY_BOTTOM - y
        rh = min(0.96, avail / max(rows, 1))
        y += max(0.0, (avail - rh * rows) / 2)
        for i, b in enumerate(bullets):
            head, desc = (b if isinstance(b, tuple) else (b, None))
            r, c = (i % rows, i // rows) if cols == 2 else (i, 0)
            x = BODY_X + c * (cw + gap)
            yy = y + r * rh
            circle(s, x + 0.06, yy + 0.13, 0.085, accent)
            iw = cw - 0.34
            budget = rh - 0.14
            hsize = autosize(head, iw, budget * (0.48 if desc else 1.0), 12,
                             bool(desc), 1.22, min_size=10)
            hh = text_height_in(head, hsize, iw, bool(desc), 1.22)
            tf = textbox(s, x + 0.24, yy, iw, budget)
            para(tf, head, size=hsize, bold=bool(desc), color=INK, line=1.22,
                 first=True)
            if desc:
                dsize = autosize(desc, iw, budget - hh - 0.06, 9.5, False, 1.3,
                                 min_size=8.0)
                para(tf, desc, size=dsize, color=MUTED, line=1.3,
                     space_before=2)
        return s

    def cards_slide(self, kicker, title, cards_data, sub=None, cols=None,
                    accent=CRIMSON, numbered=False, tinted=False):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        n = len(cards_data)
        cols = cols or (3 if n % 3 == 0 or n > 4 else (2 if n <= 4 else 3))
        rows = (n + cols - 1) // cols
        gx, gy = 0.26, 0.24
        cw = (BODY_W - gx * (cols - 1)) / cols
        avail = BODY_BOTTOM - y
        ch = (avail - gy * (rows - 1)) / rows
        pad = 0.26
        iw0 = cw - pad * 2
        need = max(
            0.30 + (0.46 if numbered else 0.0)
            + text_height_in(c[0], 12.5, iw0 - 0.19, False, 1.15, "serif")
            + 0.14
            + text_height_in(c[1], 9.8, iw0, False, 1.34) + 0.30
            for c in cards_data)
        ch = min(ch, max(need, 2.85 if rows == 1 else 1.62))
        # centre the grid, but stay near the header rather than drifting low
        y += min(0.45, max(0.0, (avail - (ch * rows + gy * (rows - 1))) / 2))

        for i, c in enumerate(cards_data):
            head, body = c[0], c[1]
            acc = c[2] if len(c) > 2 else accent
            r, cc = i // cols, i % cols
            x = BODY_X + cc * (cw + gx)
            yy = y + r * (ch + gy)
            if tinted:
                rect(s, x, yy, cw, ch, fill=acc)
                hcol, bcol = WHITE, (ON_CRIMSON if acc == CRIMSON else ON_DARK)
            else:
                card(s, x, yy, cw, ch, fill=PANEL, line=BORDER)
                vline(s, x, yy, ch, acc, 1.6)
                hcol, bcol = INK, MUTED

            contenth = ((0.46 if numbered else 0.0)
                        + text_height_in(head, 12.5, iw0 - 0.19, False, 1.15,
                                         "serif") + 0.14
                        + text_height_in(body, 9.8, iw0, False, 1.34))
            ty = yy + max(0.28, (ch - contenth) / 2)
            if numbered:
                num_badge(s, x + pad, ty, i + 1, size=0.32,
                          fill=WHITE if tinted else acc, fsize=12,
                          family="serif", color=acc if tinted else WHITE)
                ty += 0.46

            hsize = autosize(head, iw0 - 0.19, 0.78, 12.5, False, 1.15,
                             min_size=10.5, family="serif")
            hh = text_height_in(head, hsize, iw0 - 0.19, False, 1.15, "serif")
            if tinted:
                tf = textbox(s, x + pad, ty - 0.02, iw0, hh + 0.1)
                para(tf, head, size=hsize, color=hcol, family="serif",
                     line=1.15, first=True)
            else:
                marker(s, x + pad, ty + 0.075, 0.085, acc)
                tf = textbox(s, x + pad + 0.19, ty - 0.02, iw0 - 0.19, hh + 0.1)
                para(tf, head, size=hsize, color=hcol, family="serif",
                     line=1.15, first=True)

            bh = ch - (ty - yy) - hh - 0.28
            bsize = autosize(body, iw0, bh, 9.8, False, 1.34, min_size=7.5)
            tf = textbox(s, x + pad, ty + hh + 0.14, iw0, bh)
            para(tf, body, size=bsize, color=bcol, line=1.34, first=True)
        return s

    def compare_slide(self, kicker, title, left, right, sub=None,
                      lacc=CRIMSON, racc=CHARCOAL, verdict=None):
        s = self._new()
        y = self._head(s, kicker, title, sub, lacc)
        bottom = BODY_BOTTOM - (0.86 if verdict else 0.0)
        gap = 0.30
        cw = (BODY_W - gap) / 2
        h = bottom - y
        for side, x, acc in ((left, BODY_X, lacc),
                             (right, BODY_X + cw + gap, racc)):
            card(s, x, y, cw, h, fill=WHITE, line=BORDER)
            rect(s, x, y, cw, 0.34, fill=acc)
            tf = textbox(s, x + 0.24, y, cw - 0.48, 0.34,
                         anchor=MSO_ANCHOR.MIDDLE)
            para(tf, side["label"].upper(), size=8.5, color=WHITE, spacing=1.3,
                 line=1.0, first=True)
            pad = 0.28
            iw = cw - pad * 2
            hdsize = autosize(side["headline"], iw, 0.92, 14, False, 1.12,
                              min_size=11.5, family="serif")
            hh = text_height_in(side["headline"], hdsize, iw, False, 1.12,
                                "serif")
            tf = textbox(s, x + pad, y + 0.52, iw, hh + 0.1)
            para(tf, side["headline"], size=hdsize, color=INK, family="serif",
                 line=1.12, first=True)
            yy = y + 0.52 + hh + 0.20
            hline(s, x + pad, yy, iw, BORDER, 0.75)
            self._dot_rows(s, x + pad, yy + 0.20, iw, side["points"], acc,
                           avail=(y + h - 0.22) - (yy + 0.20), size=10.5)
        if verdict:
            self._band(s, BODY_BOTTOM - 0.62, verdict, label="In short",
                       fill=CRIMSON)
        return s

    def process_slide(self, kicker, title, steps, sub=None, accent=CRIMSON,
                      note=None):
        s = self._new()
        base = accent[0] if isinstance(accent, list) else accent
        y = self._head(s, kicker, title, sub, base)
        n = len(steps)
        gx = 0.20
        cw = (BODY_W - gx * (n - 1)) / n
        h = (BODY_BOTTOM - y) - (0.90 if note else 0.0)
        iw0 = cw - 0.34
        cardh = min(h - 0.72, max(
            0.22 + text_height_in(l, 11.5, iw0, True, 1.18) + 0.12
            + text_height_in(d, 9.5, iw0, False, 1.32) + 0.24
            for l, d in steps))
        cy = y + 0.34 + min(0.5, max(0.0, (h - (0.72 + cardh)) / 2))
        hline(s, BODY_X + cw * 0.5, cy - 0.008, BODY_W - cw, BORDER, 1.5)
        for i, (label, desc) in enumerate(steps):
            x = BODY_X + i * (cw + gx)
            acc = accent[i % len(accent)] if isinstance(accent, list) else accent
            num_badge(s, x + cw / 2 - 0.19, cy - 0.19, i + 1, size=0.38,
                      fill=acc, fsize=13, family="serif")
            card(s, x, cy + 0.38, cw, cardh, fill=PANEL, line=BORDER)
            hline(s, x, cy + 0.38, cw, acc, 1.4)
            lsize = autosize(label, iw0, 0.6, 11.5, True, 1.18, min_size=9.5)
            lh = text_height_in(label, lsize, iw0, True, 1.18)
            tf = textbox(s, x + 0.17, cy + 0.58, iw0, lh + 0.08)
            para(tf, label, size=lsize, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, line=1.18, first=True)
            dh = cardh - lh - 0.36
            dsize = autosize(desc, iw0, dh, 9.5, False, 1.32, min_size=7.5)
            tf = textbox(s, x + 0.17, cy + 0.58 + lh + 0.12, iw0, dh)
            para(tf, desc, size=dsize, color=MUTED, align=PP_ALIGN.CENTER,
                 line=1.32, first=True)
        if note:
            self._band(s, BODY_BOTTOM - 0.62, note, fill=base)
        return s

    def steps_slide(self, kicker, title, levels, sub=None, accent=CRIMSON):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        n = len(levels)
        gap = 0.14
        avail = BODY_BOTTOM - y
        h = min(1.02, (avail - gap * (n - 1)) / n)
        y += max(0.0, (avail - (h * n + gap * (n - 1))) / 2)
        step = (BODY_W * 0.26) / max(n - 1, 1)
        for i, lv in enumerate(levels):
            label, desc = lv[0], lv[1]
            acc = lv[2] if len(lv) > 2 else (CRIMSON if i % 2 == 0 else CHARCOAL)
            x = BODY_X + i * step
            w = BODY_W - i * step
            yy = y + (n - 1 - i) * (h + gap)
            card(s, x, yy, w, h, fill=PANEL, line=BORDER)
            vline(s, x, yy, h, acc, 1.6)
            num_badge(s, x + 0.22, yy + h / 2 - 0.16, f"L{i + 1}", size=0.32,
                      fill=acc, fsize=10, family="sans")
            lw2 = min(3.0, w - 1.15)
            lsize = autosize(label, lw2, h - 0.3, 12, True, 1.18, min_size=10)
            lh = text_height_in(label, lsize, lw2, True, 1.18)
            tf = textbox(s, x + 0.68, yy + max(0.12, (h - lh) / 2), lw2,
                         lh + 0.08)
            para(tf, label, size=lsize, bold=True, color=INK, line=1.18,
                 first=True)
            dw = w - 0.68 - lw2 - 0.5
            dsize = autosize(desc, dw, h - 0.3, 10, False, 1.3, min_size=8.0)
            dh = text_height_in(desc, dsize, dw, False, 1.3)
            tf = textbox(s, x + 0.68 + lw2 + 0.26, yy + max(0.12, (h - dh) / 2),
                         dw, dh + 0.08)
            para(tf, desc, size=dsize, color=MUTED, line=1.3, first=True)
        return s

    def stats_slide(self, kicker, title, stats, sub=None, accent=CRIMSON,
                    note=None):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        n = len(stats)
        gx = 0.26
        cw = (BODY_W - gx * (n - 1)) / n
        avail = (BODY_BOTTOM - y) - (0.86 if note else 0.0)
        h = min(2.5, avail)
        y += max(0.0, (avail - h) / 2)
        for i, (big, label, desc) in enumerate(stats):
            x = BODY_X + i * (cw + gx)
            acc = CRIMSON if i % 2 == 0 else CHARCOAL
            card(s, x, y, cw, h, fill=PANEL, line=BORDER)
            hline(s, x, y, cw, acc, 1.6)
            iw = cw - 0.52
            bigsize = autosize(big, iw, 0.86, 30, False, 1.0, min_size=15,
                               family="serif")
            bigh = text_height_in(big, bigsize, iw, False, 1.0, "serif")
            tf = textbox(s, x + 0.26, y + 0.34, iw, bigh + 0.1)
            para(tf, big, size=bigsize, color=acc, family="serif", line=1.0,
                 first=True)
            ry = y + 0.34 + bigh + 0.18
            hline(s, x + 0.26, ry, 0.5, acc, 1.2)
            lsize = autosize(label, iw, 0.56, 11.5, True, 1.18, min_size=9.5)
            lh = text_height_in(label, lsize, iw, True, 1.18)
            tf = textbox(s, x + 0.26, ry + 0.18, iw, lh + 0.08)
            para(tf, label, size=lsize, bold=True, color=INK, line=1.18,
                 first=True)
            dy = ry + 0.18 + lh + 0.10
            dh = y + h - dy - 0.24
            dsize = autosize(desc, iw, dh, 9.5, False, 1.32, min_size=7.5)
            tf = textbox(s, x + 0.26, dy, iw, dh)
            para(tf, desc, size=dsize, color=MUTED, line=1.32, first=True)
        if note:
            self._band(s, BODY_BOTTOM - 0.62, note, fill=CHARCOAL,
                       label_color=RGBColor(0x9A, 0xA4, 0xB2))
        return s

    def quadrant_slide(self, kicker, title, quads, sub=None, accent=CRIMSON):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        gx, gy = 0.26, 0.22
        cw = (BODY_W - gx) / 2
        ch = (BODY_BOTTOM - y - gy) / 2
        for i, q in enumerate(quads):
            head, pts = q[0], q[1]
            acc = q[2] if len(q) > 2 else (CRIMSON if i % 2 == 0 else CHARCOAL)
            x = BODY_X + (i % 2) * (cw + gx)
            yy = y + (i // 2) * (ch + gy)
            card(s, x, yy, cw, ch, fill=PANEL, line=BORDER)
            vline(s, x, yy, ch, acc, 1.6)
            hh = self._panel_head(s, x + 0.28, yy + 0.22, head, cw - 0.56, acc,
                                  size=12.5)
            yy2 = yy + 0.22 + hh + 0.16
            hline(s, x + 0.28, yy2 - 0.08, cw - 0.56, BORDER, 0.75)
            self._dot_rows(s, x + 0.28, yy2 + 0.04, cw - 0.56, pts, acc,
                           avail=(yy + ch - 0.18) - (yy2 + 0.04), size=10)
        return s

    def table_slide(self, kicker, title, headers, rows, sub=None,
                    accent=CRIMSON, widths=None, note=None):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        ncol, nrow = len(headers), len(rows) + 1
        avail = (BODY_BOTTOM - y) - (0.86 if note else 0.0)
        h = min(avail, 0.46 + 0.56 * len(rows))
        y += max(0.0, (avail - h) / 2)
        gt = s.shapes.add_table(nrow, ncol, _i(BODY_X), _i(y), _i(BODY_W), _i(h))
        tbl = gt.table
        for tag in ("bandRow", "bandCol", "firstCol", "lastRow", "lastCol"):
            tbl._tbl.tblPr.set(tag, "0")
        tbl._tbl.tblPr.set("firstRow", "1")
        if widths:
            total = sum(widths)
            for i, w in enumerate(widths):
                tbl.columns[i].width = Emu(int(Inches(BODY_W) * w / total))
        colw = [tbl.columns[i].width / 914400.0 for i in range(ncol)]
        rowh = max(0.40, (h - 0.46) / len(rows))
        tbl.rows[0].height = Inches(0.46)
        for r in range(1, nrow):
            tbl.rows[r].height = Inches(rowh)

        body_size = 10.0
        while body_size > 7.5:
            worst = max(text_height_in(v, body_size, max(colw[c] - 0.34, 0.5),
                                       c == 0, 1.26)
                        for row in rows for c, v in enumerate(row))
            if worst <= rowh - 0.16:
                break
            body_size -= 0.5
        head_size = 10.0
        while head_size > 8.0 and max(
                text_height_in(hh, head_size, max(colw[c] - 0.34, 0.5), True,
                               1.15)
                for c, hh in enumerate(headers)) > 0.28:
            head_size -= 0.5

        for c, htxt in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = accent
            cell.margin_left = cell.margin_right = Inches(0.16)
            cell.margin_top = cell.margin_bottom = Inches(0.08)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            para(cell.text_frame, htxt.upper(), size=head_size, color=WHITE,
                 spacing=1.1, line=1.15, first=True)
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else PANEL
                cell.margin_left = cell.margin_right = Inches(0.16)
                cell.margin_top = cell.margin_bottom = Inches(0.08)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.word_wrap = True
                para(cell.text_frame, val, size=body_size, bold=(c == 0),
                     color=INK if c == 0 else MUTED, line=1.26, first=True)
        if note:
            self._band(s, BODY_BOTTOM - 0.62, note, fill=accent)
        return s

    def split_slide(self, kicker, title, left_head, left_points, right_cards,
                    sub=None, accent=CRIMSON):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        lw = BODY_W * 0.40
        rw = BODY_W - lw - 0.40
        h = BODY_BOTTOM - y
        lhsize = autosize(left_head, lw, 1.5, 16, False, 1.14, min_size=13,
                          family="serif")
        lhh = text_height_in(left_head, lhsize, lw, False, 1.14, "serif")
        tf = textbox(s, BODY_X, y + 0.02, lw, lhh + 0.14)
        para(tf, left_head, size=lhsize, color=INK, family="serif", line=1.14,
             first=True)
        yy = y + 0.04 + lhh + 0.22
        hline(s, BODY_X, yy, 0.62, accent, 1.6)
        yy += 0.28
        pw = lw - 0.26
        avail = (y + h) - yy
        gapp = 0.16
        psize = 10.5
        while psize > 8.5 and (sum(text_height_in(p, psize, pw, False, 1.3)
                                   for p in left_points)
                               + gapp * (len(left_points) - 1)) > avail:
            psize -= 0.5
        for p in left_points:
            ph = text_height_in(p, psize, pw, False, 1.3)
            circle(s, BODY_X + 0.055, yy + 0.115, 0.075, accent)
            tfp = textbox(s, BODY_X + 0.22, yy, pw, ph + 0.08)
            para(tfp, p, size=psize, color=INK, line=1.3, first=True)
            yy += ph + gapp

        x = BODY_X + lw + 0.40
        n = len(right_cards)
        gy = 0.18
        ch = (h - gy * (n - 1)) / n
        for i, (hd, bd) in enumerate(right_cards):
            acc = CRIMSON if i % 2 == 0 else CHARCOAL
            yy2 = y + i * (ch + gy)
            card(s, x, yy2, rw, ch, fill=PANEL, line=BORDER)
            vline(s, x, yy2, ch, acc, 1.6)
            iw = rw - 0.54
            hsize = autosize(hd, iw - 0.19, 0.5, 12.5, False, 1.14,
                             min_size=10.5, family="serif")
            hh = text_height_in(hd, hsize, iw - 0.19, False, 1.14, "serif")
            bsize = autosize(bd, iw, ch - hh - 0.44, 9.8, False, 1.32,
                             min_size=7.5)
            bh = text_height_in(bd, bsize, iw, False, 1.32)
            top = yy2 + max(0.16, (ch - hh - bh - 0.12) / 2)
            marker(s, x + 0.28, top + 0.07, 0.085, acc)
            tf = textbox(s, x + 0.47, top - 0.02, iw - 0.19, hh + 0.1)
            para(tf, hd, size=hsize, color=INK, family="serif", line=1.14,
                 first=True)
            tf = textbox(s, x + 0.28, top + hh + 0.12, iw, bh + 0.1)
            para(tf, bd, size=bsize, color=MUTED, line=1.32, first=True)
        return s

    def quote_slide(self, statement, attribution=None, kicker=None,
                    accent=CRIMSON):
        s = self._new(footer=False, bg=None)
        gradient_bg(s, CRIMSON, CRIMSON_DK, 315.0)
        ghost(s, -1.8, 3.8, 5.6, 5.6, "000000", 2600)
        ghost(s, 10.4, -1.6, 5.6, 5.6, "FFFFFF", 2600)
        if self.logo:
            place_logo(s, self.logo, 0, 0.62, 0.7, on_dark=True, max_w=2.0,
                       right_edge=ML + CW - 0.06)
        if kicker:
            w = chip_width(kicker, 8.5, False, 1.3, 0.44)
            chip(s, SW / 2 - w / 2, 2.30, kicker, fill=CRIMSON_MID, color=WHITE,
                 size=8.5, h=0.30, pad=0.44)
        qw = SW - 3.4
        qsize = autosize(statement, qw, 2.1, 30, False, 1.18, min_size=20,
                         family="serif")
        qh = text_height_in(statement, qsize, qw, False, 1.18, "serif")
        tf = textbox(s, 1.7, 3.00 + (2.1 - qh) / 2, qw, qh + 0.2)
        para(tf, statement, size=qsize, color=WHITE, family="serif", line=1.18,
             align=PP_ALIGN.CENTER, first=True)
        hline(s, SW / 2 - 0.42, 5.34, 0.84, ON_CRIMSON_LT, 1.4)
        if attribution:
            tf = textbox(s, 2.3, 5.62, SW - 4.6, 0.6)
            para(tf, attribution, size=11.5, color=ON_CRIMSON,
                 align=PP_ALIGN.CENTER, line=1.34, first=True)
        caption(s, ML, FOOTER_Y,
                f"{self.PRESENTER}  ·  Register No. {self.REGISTER}", w=5.0,
                color=ON_CRIMSON_LT)
        return s

    def takeaways_slide(self, points, title="Key takeaways", kicker="Recap",
                        sub=None, accent=CRIMSON):
        s = self._new()
        y = self._head(s, kicker, title, sub, accent)
        n = len(points)
        cols = 2 if n > 5 else 1
        rows = (n + cols - 1) // cols
        gx, gy = 0.28, 0.20
        cw = (BODY_W - gx) / 2 if cols == 2 else BODY_W
        avail = BODY_BOTTOM - y
        h = min(1.5, (avail - gy * (rows - 1)) / rows)
        y += max(0.0, (avail - (h * rows + gy * (rows - 1))) / 2)
        for i, p in enumerate(points):
            head, desc = (p if isinstance(p, tuple) else (p, None))
            r, c = (i % rows, i // rows) if cols == 2 else (i, 0)
            x = BODY_X + c * (cw + gx)
            yy = y + r * (h + gy)
            acc = CRIMSON if i % 2 == 0 else CHARCOAL
            card(s, x, yy, cw, h, fill=PANEL, line=BORDER)
            vline(s, x, yy, h, acc, 1.6)
            num_badge(s, x + 0.26, yy + h / 2 - 0.16, i + 1, size=0.32,
                      fill=acc, fsize=12, family="serif")
            iw = cw - 1.06
            budget = h - 0.24
            hsize = autosize(head, iw, budget * (0.54 if desc else 1.0), 12,
                             True, 1.18, min_size=10)
            hh = text_height_in(head, hsize, iw, True, 1.18)
            dsize = dh = 0
            if desc:
                dsize = autosize(desc, iw, budget - hh - 0.06, 9.8, False, 1.32,
                                 min_size=8.0)
                dh = text_height_in(desc, dsize, iw, False, 1.32) + 0.04
            top = yy + max(0.12, (h - hh - dh) / 2)
            tf = textbox(s, x + 0.74, top, iw, hh + dh + 0.1)
            para(tf, head, size=hsize, bold=True, color=INK, line=1.18,
                 first=True)
            if desc:
                para(tf, desc, size=dsize, color=MUTED, line=1.32,
                     space_before=2)
        return s

    def closing_slide(self, title="Thank you", subtitle=None, questions=None):
        s = self._new(footer=False, bg=None)
        gradient_bg(s, CHARCOAL, CHARCOAL_DK, 315.0)
        ghost(s, 9.4, -1.6, 6.0, 6.0, "FFFFFF", 2400)
        ghost(s, 11.4, 3.8, 4.0, 4.0, "8B0000", 9000)
        if self.logo:
            place_logo(s, self.logo, 0, 0.62, 0.8, on_dark=True, max_w=2.3,
                       right_edge=ML + CW - 0.06)
        x = ML + 0.32
        chip(s, x, 1.74, "Conclusion", fill=CRIMSON, color=WHITE, size=8.5,
             h=0.30, pad=0.44)
        tsize = autosize(title, 8.6, 1.2, 40, False, 1.02, min_size=30,
                         family="serif")
        th = text_height_in(title, tsize, 8.6, False, 1.02, "serif")
        tf = textbox(s, x, 2.22, 8.6, th + 0.16)
        para(tf, title, size=tsize, color=WHITE, family="serif", line=1.02,
             first=True)
        y = 2.22 + th + 0.24
        hline(s, x, y, 0.86, CRIMSON, 1.6)
        if subtitle:
            ssize = autosize(subtitle, 8.4, 0.9, 12.5, False, 1.38,
                             min_size=10.5)
            sh = text_height_in(subtitle, ssize, 8.4, False, 1.38)
            tf = textbox(s, x, y + 0.28, 8.4, sh + 0.1)
            para(tf, subtitle, size=ssize, color=ON_DARK, line=1.38, first=True)
            y = y + 0.28 + sh
        if questions:
            qy = max(y + 0.36, 4.30)
            n = len(questions)
            gx = 0.24
            qw = (CW - gx * (n - 1)) / n
            qh = min(1.30, 5.98 - qy)
            for i, q in enumerate(questions):
                qx = ML + i * (qw + gx)
                rect(s, qx, qy, qw, qh, fill=RGBColor(0x2B, 0x35, 0x44),
                     line=None)
                vline(s, qx, qy, qh, CRIMSON, 1.6)
                iw = qw - 0.46
                qs = autosize(q, iw, qh - 0.54, 10, False, 1.32, min_size=8.0)
                eyebrow(s, qx + 0.24, qy + 0.20, f"Question {i + 1}",
                        color=RGBColor(0x8E, 0x99, 0xA8), size=7.5, w=iw)
                tf = textbox(s, qx + 0.24, qy + 0.46, iw, qh - 0.62)
                para(tf, q, size=qs, color=RGBColor(0xD8, 0xDE, 0xE6),
                     line=1.32, first=True)
        # presenter sign-off bar
        rect(s, ML, 6.30, CW, 0.62, fill=CRIMSON)
        tf = textbox(s, ML + 0.30, 6.30, CW - 0.6, 0.62,
                     anchor=MSO_ANCHOR.MIDDLE)
        rich(tf, [("Presented by   ", {"size": 8.0, "color": ON_CRIMSON_LT,
                                       "spacing": 1.3, "caps": True}),
                  (self.PRESENTER, {"size": 13, "color": WHITE,
                                    "family": "serif"}),
                  ("      Register No.   ", {"size": 8.0,
                                             "color": ON_CRIMSON_LT,
                                             "spacing": 1.3, "caps": True}),
                  (self.REGISTER, {"size": 13, "color": WHITE,
                                   "family": "serif"})],
             first=True)
        return s


# --------------------------------------------------------------------------
def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text
    return slide
